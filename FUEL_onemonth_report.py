# Importing necessary libraries
import pandas as pd
import numpy as np
import math
from pptx import Presentation
from pptx.util import Inches
from datetime import datetime

def load_csv_to_dataframe(file_path):
    """
    Load a CSV file into a pandas DataFrame.

    Args:
    file_path (str): The path to the CSV file.

    Returns:
    pd.DataFrame: The loaded DataFrame.
    """
    df = pd.read_csv(file_path, sep='\t')
    return df

def process_date_column(df):
    """
    Convert 'Date' column to datetime and sort the DataFrame by 'Date' in descending order.

    Args:
    df (pd.DataFrame): The DataFrame with a 'Date' column.

    Returns:
    pd.DataFrame: The sorted DataFrame.
    """
    df['Date'] = pd.to_datetime(df['Date'], format='%Y-%m-%d')
    df_sorted = df.sort_values('Date', ascending=False)
    return df_sorted

def pre_process_dataframe(df):
    """
    The function preprocesses the input DataFrame by replacing any values less than or equal to 0 
    (or equal to -1) in the specified columns with NaN.

    Args:
    df: pd.DataFrame - input DataFrame to be preprocessed

    Returns:
    df: pd.DataFrame - preprocessed DataFrame
    """
    columns_to_preprocess = ['HabitsSteps', 'MyHabitsNutritionEvolution', 'HabitsSleepQuantity', 
                             'HabitsPercentageSleepRestful', 'MyHabitsSleepEvolution', 'BodyRestingHeartRate', 
                             'BodyHeartRateVariability', 'BodyMassIndex', 'BodyTotalBodyWater', 
                             'BodyWaistCircumference', 'BodyPercentageBodyFat', 'BehaviorPerception', 
                             'BehaviorExpression', 'BehaviorOrientation', 'BehaviorConnection', 'BehaviorsApplication']

    for column in columns_to_preprocess:
        df[column] = df[column].where(df[column] > 0, np.nan)
        
    df[df == -1] = np.nan

    return df

def calculate_week_dates(df):
    """
    Calculate the current week dates and last week dates based on the maximum date in the DataFrame.
    It also separates the dates into weekends and weekdays.

    Args:
    df (pd.DataFrame): The DataFrame with a 'Date' column.

    Returns:
    tuple: A tuple containing four lists: the current week dates, the last week dates, current weekend dates, and current weekday dates.
    """
    # Calculate current week dates
    max_date = df['Date'].max()
    curr_date_week = [max_date - pd.Timedelta(days=29), max_date]

    # Calculate last week dates
    last_date_week = [(date - pd.Timedelta(days=30)).strftime('%Y-%m-%d') for date in curr_date_week]

    # Format dates as strings in the "YYYY-MM-DD" format
    curr_date_week = [date.strftime('%Y-%m-%d') for date in curr_date_week]

    # Create a date range for current week
    curr_week = pd.date_range(start=curr_date_week[0], end=curr_date_week[1])

    # Check if each date in curr_week is a weekend
    is_weekend = curr_week.to_series().dt.dayofweek >= 5

    # Create arrays for weekend and weekday dates
    curr_weekend_arr = curr_week[is_weekend].strftime('%Y-%m-%d')
    curr_week_arr = curr_week[~is_weekend].strftime('%Y-%m-%d')

    return curr_date_week, last_date_week, curr_weekend_arr.tolist(), curr_week_arr.tolist()

def analyze_participant_data(df, currDateWeek, currWeekendArr, currWeekArr, lastDateWeek):
    """
    Analyze participant data for specific health indicators, calculate the average steps
    and total activity time, provide recommendations to improve their activity level,
    and analyze nutrition data.
    """
    mPartVecIds = df['ParticipantProgramId'].unique()
    for mpart_id in mPartVecIds:
    
        print(f"###########################Current mpart_id##########################: {mpart_id}")
        try:
    
            if mpart_id != 222:
                
                # Analyze participant sleep data
                print("************SLEEP - START************")
                act1Sleep, mSleepActions, mSleepDispersionGreen, mSleepDispersionRed, act2Sleep, mSleepQualityRed, mSleepQualityGreen, act3Sleep, avgCurrSleepMins, mSleepTimeRed, mSleepTimeGreen, mSleepTime = analyze_sleep_data(df, mpart_id, currDateWeek, currWeekendArr, currWeekArr, lastDateWeek)
                # Print the variables
                print("act1Sleep: ", act1Sleep)
                print("mSleepActions:", mSleepActions)
                print("mSleepDispersionGreen: ", mSleepDispersionGreen)
                print("mSleepDispersionRed: ", mSleepDispersionRed)
                print("act2Sleep: ", act2Sleep)
                print("mSleepQualityRed: ", mSleepQualityRed)
                print("mSleepQualityGreen: ", mSleepQualityGreen)
                print("act3Sleep: ", act3Sleep)
                print("avgCurrSleepMins: ", avgCurrSleepMins)
                print("mSleepTimeRed: ", mSleepTimeRed)
                print("mSleepTimeGreen: ", mSleepTimeGreen)
                print("mSleepTime: ", mSleepTime)
                print("************SLEEP - END************")

                # Calculate Participant Steps
                print("************STEPS - START************")
                mActivitytimeActions, mActivitytime, mSteps, mActivitytimeRed, mActivitytimeGreen, mStepsLastRed, mStepsLastGreen = calculate_steps(df, mpart_id, currDateWeek, currWeekendArr, currWeekArr, lastDateWeek) 
                print("m_activity_time_actions:", mActivitytimeActions)
                print("m_activity_time:", mActivitytime)
                print("m_steps:", mSteps)
                print("mActivitytimeRed:", mActivitytimeRed)
                print("mActivitytimeGreen:", mActivitytimeGreen)
                print("mStepsLastRed:", mStepsLastRed)
                print("mStepsLastGreen:", mStepsLastGreen)
                print("************STEPS - END************")

                # Processes nutrition data
                print("************Processes nutrition data - START************")
                mNutritionscore, mNutritionscoreRed, mNutritionscoreGreen, avgCurrMeals, mDailymeals, mDailymealsGreen, mDailymealsRed, mRedNutrition1, mRedNutrition2, mRedNutrition3, mRedNutrition4, mRedNutrition5, mRedLabelNutrition1, mRedLabelNutrition2, mRedLabelNutrition3, mRedLabelNutrition4, mRedLabelNutrition5, mNutrition1, mNutrition2, mNutrition3, mNutrition4, mNutrition5, mGreenNutrition1, mGreenNutrition2, mGreenNutrition3, mGreenNutrition4, mGreenNutrition5, fracValueComparison = process_nutrition_data(df, mpart_id, currDateWeek, lastDateWeek)
                # Print the results
                print("mNutritionscore: ", mNutritionscore)
                print("avgCurrMeals: ", avgCurrMeals)
                print("mDailymeals: ", mDailymeals)
                print("mDailymealsGreen: ", mDailymealsGreen)
                print("mDailymealsRed: ", mDailymealsRed)
                print("mNutritionscoreRed: ", mNutritionscoreRed)
                print("mNutritionscoreGreen: ", mNutritionscoreGreen)
                print("mRedNutrition1 =", mRedNutrition1)
                print("mRedLabelNutrition1 =", mRedLabelNutrition1)
                print("mRedNutrition2 =", mRedNutrition2)
                print("mRedLabelNutrition2 =", mRedLabelNutrition2)
                print("mRedNutrition3 =", mRedNutrition3)
                print("mRedLabelNutrition3 =", mRedLabelNutrition3)
                print("mRedNutrition4 =", mRedNutrition4)
                print("mRedLabelNutrition4 =", mRedLabelNutrition4)
                print("mRedNutrition5 =", mRedNutrition5)
                print("mRedLabelNutrition5 =", mRedLabelNutrition5)
                print("mGreenNutrition1 =", mGreenNutrition1)
                print("mNutrition1 =", mNutrition1)
                print("mGreenNutrition2 =", mGreenNutrition2)
                print("mNutrition2 =", mNutrition2)
                print("mGreenNutrition3 =", mGreenNutrition3)
                print("mNutrition3 =", mNutrition3)
                print("mGreenNutrition4 =", mGreenNutrition4)
                print("mNutrition4 =", mNutrition4)
                print("mGreenNutrition5 =", mGreenNutrition5)
                print("mNutrition5 =", mNutrition5)
                print("fracValueComparison: ", fracValueComparison)
                print("************Processes nutrition data - END************")

                # Calculate nutrition actions
                action1Nutri = calculate_nutrition_actions(df, mpart_id, currWeekendArr, currWeekArr)

                # Calculate meal actions
                action2Nutri = calculate_meal_actions(avgCurrMeals)

                # Calculate fruit actions
                action3Nutri = calculate_fruit_actions(df, mpart_id, currDateWeek, fracValueComparison)

                # Calculate sugar actions
                action4Nutri = calculate_sugar_actions(df, mpart_id, currDateWeek, fracValueComparison)

                # Calculate whole grain actions
                action5Nutri = calculate_whole_grain_actions(df, mpart_id, currDateWeek, fracValueComparison)

                # Calculate protein actions
                action6Nutri = calculate_protein_actions(df, mpart_id, currDateWeek, fracValueComparison)

                # Calculate water actions
                action7Nutri = calculate_water_actions(df, mpart_id, currDateWeek, fracValueComparison)

                # Calculate vegetable actions
                action8Nutri = calculate_vegetable_actions(df, mpart_id, currDateWeek)

                # Calculate nutrition actions
                print("************Consolidate nutrition actions - START************")
                mActionsNutrition = consolidate_nutrition_actions(action1Nutri, action2Nutri, action3Nutri, action4Nutri, action5Nutri, action6Nutri, action7Nutri, action8Nutri)
                print(mActionsNutrition)
                print("************Consolidate nutrition actions - END************")

                #Generate the pptx presentation
                print("************Generate Presentation - START************")
                # Convert strings to datetime objects
                date1 = datetime.strptime(currDateWeek[0], '%Y-%m-%d')
                date2 = datetime.strptime(currDateWeek[1], '%Y-%m-%d')
                # Format dates and create date range string
                DateValue = "{}/{} - {}/{}".format(date1.day, date1.month, date2.day, date2.month)

                data = {
                    "ReportDate": DateValue,
                    "SleepTime": mSleepTime,
                    "SleepTimeRed": mSleepTimeRed,
                    "SleepTimeGreen": mSleepTimeGreen,
                    "SleepDispersionGreen": mSleepDispersionGreen,
                    "SleepDispersionRed": mSleepDispersionRed,
                    "SleepQualityRed": mSleepQualityRed,
                    "SleepQualityGreen": mSleepQualityGreen,
                    "SleepActions": mSleepActions,
                    "Steps": mSteps,
                    "StepsLastGreen": mStepsLastGreen,
                    "StepsLastRed": mStepsLastRed,
                    "Activitytime": mActivitytime,
                    "ActivitytimeGreen": mActivitytimeGreen,
                    "ActivitytimeRed": mActivitytimeRed,
                    "ActivitytimeActions": mActivitytimeActions,
                    "Nutritionscore": mNutritionscore["MyHabitsNutritionEvolution"],
                    "NutritionscoreRed": mNutritionscoreRed["MyHabitsNutritionEvolution"],
                    "NutritionscoreGreen": mNutritionscoreGreen["MyHabitsNutritionEvolution"],
                    "Dailymeals": mDailymeals["NutritionDailyMeals"],
                    "DailymealsGreen": mDailymealsGreen["NutritionDailyMeals"],
                    "DailymealsRed": mDailymealsRed["NutritionDailyMeals"],
                    "GreenNutrition1": mGreenNutrition1,
                    "Nutrition1": mNutrition1,
                    "GreenNutrition2": mGreenNutrition2,
                    "Nutrition2": mNutrition2,
                    "GreenNutrition3": mGreenNutrition3,
                    "Nutrition3": mNutrition3,
                    "GreenNutrition4": mGreenNutrition4,
                    "Nutrition4": mNutrition4,
                    "GreenNutrition5": mGreenNutrition5,
                    "Nutrition5": mNutrition5,
                    "RedNutrition1": mRedNutrition1,
                    "RedLabelNutrition1": mRedLabelNutrition1,
                    "RedNutrition2": mRedNutrition2,
                    "RedLabelNutrition2": mRedLabelNutrition2,
                    "RedNutrition3": mRedNutrition3,
                    "RedLabelNutrition3": mRedLabelNutrition3,
                    "RedNutrition4": mRedNutrition4,
                    "RedLabelNutrition4": mRedLabelNutrition4,
                    "RedNutrition5": mRedNutrition5,
                    "RedLabelNutrition5": mRedLabelNutrition5,
                    "ActionsNutrition": mActionsNutrition
                }
                create_presentation(data, mpart_id)
                print("************Generate Presentation - END************")

        except Exception as e:
            print(f"Exception occurred for mpart_id {mpart_id}: {e}")

    return mpart_id

def create_presentation(data, participant_id):
    print("INSIDE create_presentation")

     # Load the PowerPoint template
    presentation = Presentation('weekReport.pptx')
    slide = presentation.slides[0]

    #Placeholder index to data key mapping
    idx_to_data_key = {
                15: 'ReportDate',
                66: 'RedLabelNutrition5',
                67: 'RedNutrition5',
                68: 'Nutrition5',
                69: 'GreenNutrition5',
                45: 'RedLabelNutrition1',
                63: 'RedNutrition1',
                44: 'Nutrition1',
                65: 'GreenNutrition1',
                27: 'NutritionscoreRed',
                28: 'StepsLastGreen',
                29: 'StepsLastRed',
                30: 'ActivitytimeGreen',
                31: 'ActivitytimeRed',
                37: 'SleepQualityGreen',
                33: 'SleepDispersionRed',
                34: 'SleepTimeGreen',
                35: 'SleepTimeRed',
                20: 'Activitytime',
                21: 'Steps',
                18: 'Dailymeals',
                19: 'Nutritionscore',
                13: 'SleepTime',
                26: 'NutritionscoreGreen',
                32: 'SleepDispersionGreen',
                36: 'SleepQualityRed',
                42: 'DailymealsGreen',
                43: 'DailymealsRed',
                70: 'RedLabelNutrition2',
                71: 'RedNutrition2',
                72: 'Nutrition2',
                73: 'GreenNutrition2',
                74: 'RedLabelNutrition3',
                75: 'RedNutrition3',
                76: 'Nutrition3',
                77: 'GreenNutrition3',
                78: 'RedLabelNutrition4',
                79: 'RedNutrition4',
                80: 'Nutrition4',
                81: 'GreenNutrition4',
                82: 'SleepActions',
                83: 'ActivitytimeActions',
                84: 'ActionsNutrition'
            }

    # Process each shape in the slide
    for shape in slide.shapes:
        if shape.is_placeholder:
            placeholder_idx = shape.placeholder_format.idx

            if placeholder_idx in idx_to_data_key:
                # Clear existing text in the text frame
                text_frame = shape.text_frame
                text_frame.clear()
                
                # Manage initial paragraph to avoid an empty first line
                if text_frame.paragraphs:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()

                # Handle bullet point formatting
                if placeholder_idx in [82, 83, 84]:  # IDs for placeholders that need bullet points
                    actions_text = data[idx_to_data_key[placeholder_idx]].split('\n')
                    first_action = True
                    for action in actions_text:
                        if not first_action:
                            p = text_frame.add_paragraph()  # add new paragraph after the first
                        p.text = action
                        p.level = 0  # Set bullet point level (0 for top-level)
                        p.space_before = 0
                        p.space_after = 0
                        first_action = False
                else:
                    # Regular text
                    p.text = str(data[idx_to_data_key[placeholder_idx]])

    # Save the presentation
    filenamepptx = f"{participant_id}_Relatório Semanal.pptx"
    presentation.save(filenamepptx)

def analyze_sleep_data(df, mpartID, currDateWeek, currWeekendArr, currWeekArr, lastDateWeek):
    """
    Analyze participant sleep data for specific health indicators, calculate the average sleep
    quantity and quality, and provide recommendations to improve their sleep habits.
    """
    currWeekSleep = df[(df['ParticipantProgramId'] == mpartID) & 
                       (df['DataType'] == "I") & 
                       (df['Date'] >= currDateWeek[0]) & 
                       (df['Date'] <= currDateWeek[1])].filter(like="HabitsSleepQuantity")
    
    lastWeekSleep = df[(df['ParticipantProgramId'] == mpartID) & 
                       (df['DataType'] == "I") & 
                       (df['Date'] >= lastDateWeek[0]) & 
                       (df['Date'] <= lastDateWeek[1])].filter(like="HabitsSleepQuantity")
    
    currWeekSleepQlt = df[(df['ParticipantProgramId'] == mpartID) & 
                          (df['DataType'] == "G") & 
                          (df['Date'] >= currDateWeek[0]) & 
                          (df['Date'] <= currDateWeek[1])].filter(like="HabitsPercentageSleepRestful")
    
    SleepQltAvg = currWeekSleepQlt.mean(numeric_only=True, skipna=True)
    currWeekSleepMinsSd = currWeekSleep.apply(lambda x: np.std(x*60, ddof=1))

    act2Sleep = ""
    mSleepDispersionRed = " "
    mSleepDispersionGreen = " "

    # Check if currWeekSleepMinsSd or SleepQltAvg is NaN
    if currWeekSleepMinsSd.isna().any().item() or SleepQltAvg.isna().any().item():
        act2Sleep = "melhorar a inserção de dados do sono"
    else:
        currWeekSleepMinsSd = currWeekSleepMinsSd.item()
        if 0 <= currWeekSleepMinsSd < 30:
            mSleepDispersionGreen = "reduzida"
            mSleepDispersionRed = " "
            act2Sleep = "manter rotina de sono"
        elif 30 <= currWeekSleepMinsSd < 60:
            mSleepDispersionGreen = "moderada"
            mSleepDispersionRed = " "
            act2Sleep = "manter rotina de sono"
        elif 60 <= currWeekSleepMinsSd < 90:
            mSleepDispersionGreen = " "
            mSleepDispersionRed = "elevada"
            act2Sleep = "estabilizar horários de sono"
        elif 90 <= currWeekSleepMinsSd < 120:
            mSleepDispersionGreen = " "
            mSleepDispersionRed = "muito elevada"
            act2Sleep = "estabilizar horários de sono"
        elif currWeekSleepMinsSd >= 120:
            mSleepDispersionGreen = " "
            mSleepDispersionRed = "extremamente elevada"
            act2Sleep = "estabilizar horários de sono"
    
    act3Sleep = ""
    SleepQltAvg = SleepQltAvg.item()
    if 0 <= SleepQltAvg < 10:
        mSleepQualityRed = "muito má"
        mSleepQualityGreen = " "
        act3Sleep = "melhorar higiene de sono"
    elif 10 <= SleepQltAvg < 30:
        mSleepQualityRed = "má"
        mSleepQualityGreen = " "
        act3Sleep = "melhorar higiene de sono"
    elif 30 <= SleepQltAvg < 50:
        mSleepQualityRed = "fraca"
        mSleepQualityGreen = " "
        act3Sleep = "melhorar higiene de sono"
    elif 50 <= SleepQltAvg < 60:
        mSleepQualityRed = "razoável"
        mSleepQualityGreen = " "
        act3Sleep = "melhorar higiene de sono"
    elif 60 <= SleepQltAvg < 70:
        mSleepQualityRed = " "
        mSleepQualityGreen = "boa"
        act3Sleep = "melhorar higiene de sono"
    elif 70 <= SleepQltAvg < 80:
        mSleepQualityRed = " "
        mSleepQualityGreen = "muito boa"
        act3Sleep = "melhorar higiene de sono"
    elif 80 <= SleepQltAvg < 90:
        mSleepQualityRed = " "
        mSleepQualityGreen = "excelente"
        act3Sleep = ""
    elif SleepQltAvg >= 90:
        mSleepQualityRed = " "
        mSleepQualityGreen = "de atleta"
        act3Sleep = ""

    avgCurrSleep = currWeekSleep.astype(float).mean(skipna=True)
    avgLastSleep = lastWeekSleep.astype(float).mean(skipna=True)

    diffSleep = avgCurrSleep - avgLastSleep

    hours = int(np.floor(avgCurrSleep.iloc[0]))
    mins = int(np.floor((((avgCurrSleep.iloc[0] * 10) % 10 / 10) * 60) + 0.5))

    if hours < 10:
        hours = f"{hours:02d}"
    else:
        hours = f"{hours:.0f}"
        
    if mins < 10:
        mins = f"{mins:02d}"
    else:
        mins = f"{mins:.0f}"
        
    mSleepTime = f"{hours}:{mins}"
        
    if pd.notna(diffSleep).any():
        sleepABS = abs(diffSleep)
        hours = int(np.floor(sleepABS.iloc[0]))
        mins = int(np.floor((((sleepABS.iloc[0] * 10) % 10 / 10) * 60) + 0.5))

        if hours < 10:
            hours = f"{abs(hours):02d}"
        else:
            hours = f"{abs(hours):.0f}"

        if mins < 10:
            mins = f"{mins:02d}"
        else:
            mins = f"{mins:.0f}"

        if (diffSleep >= 0).any():
            mSleepTimeGreen = f"+{hours}:{mins}"
            mSleepTimeRed = " "
        else:
            mSleepTimeRed = f"-{hours}:{mins}"
            mSleepTimeGreen = " "
    else:
        mSleepTimeGreen = "No Data Available"
        mSleepTimeRed = "No Data Available"

    avgCurrSleepMins = avgCurrSleep * 60

    avgWeekendSleep = 0
    countSleepUnits = 0
    for i in range(len(currWeekendArr)):
        weekendSleep = df[(df['ParticipantProgramId'] == mpartID) & 
                        (df['DataType'] == "I") & 
                        (df['Date'] == currWeekendArr[i])].filter(like="HabitsSleepQuantity")
        if weekendSleep.notna().values.any():
            avgWeekendSleep += weekendSleep.astype(float).sum().item()  # Ensure it's scalar
            countSleepUnits += 1

    if countSleepUnits > 0:
        avgWeekendSleep = avgWeekendSleep / countSleepUnits
    else:
        avgWeekendSleep = np.nan  # Handle case where no data is available

    avgWeekSleep = 0
    countSleepUnits = 0
    for i in range(len(currWeekArr)):
        weekSleep = df[(df['ParticipantProgramId'] == mpartID) & 
                    (df['DataType'] == "I") & 
                    (df['Date'] == currWeekArr[i])].filter(like="HabitsSleepQuantity")
        if weekSleep.notna().values.any():
            avgWeekSleep += weekSleep.astype(float).sum().item()  # Ensure it's scalar
            countSleepUnits += 1

    if countSleepUnits > 0:
        avgWeekSleep = avgWeekSleep / countSleepUnits
    else:
        avgWeekSleep = np.nan  # Handle case where no data is available

    avgWeekSleepMins = avgWeekSleep * 60 if isinstance(avgWeekSleep, float) else np.nan
    avgWeekendSleepMins = avgWeekendSleep * 60 if isinstance(avgWeekendSleep, float) else np.nan

    # diffWeekWeekend calculation assuming both are floats or NaN
    diffWeekWeekend = avgWeekendSleepMins - avgWeekSleepMins if pd.notna(avgWeekendSleepMins) and pd.notna(avgWeekSleepMins) else np.nan

    # No need for .mean() if avgWeekSleepMins is already a float
    avgWeekSleepMins_value = avgWeekSleepMins

    act1Sleep = ""

    # First, ensure that diffWeekWeekend is not NaN and handle both Series and scalar cases.
    if isinstance(diffWeekWeekend, pd.Series):
        is_not_na = pd.notna(diffWeekWeekend).all()  # Use .all() when it's a Series
    else:
        is_not_na = pd.notna(diffWeekWeekend)  # Direct check if it's a scalar

    if is_not_na:
        if avgWeekSleepMins_value < 480:
            # Check conditions based on whether diffWeekWeekend is a Series or scalar
            if isinstance(diffWeekWeekend, pd.Series):
                condition1 = ((diffWeekWeekend > -60) & (diffWeekWeekend < 60)).any()
                condition2 = (diffWeekWeekend >= 60).any()
                condition3 = (diffWeekWeekend < -60).any()
            else:
                condition1 = -60 < diffWeekWeekend < 60
                condition2 = diffWeekWeekend >= 60
                condition3 = diffWeekWeekend < -60

            if condition1:
                act1Sleep = "dormir mais tempo"
            elif condition2:
                act1Sleep = "dormir mais durante a semana"
            elif condition3:
                act1Sleep = "dormir mais nos fins de semana"
        else:
            if isinstance(diffWeekWeekend, pd.Series):
                condition1 = ((diffWeekWeekend > -60) & (diffWeekWeekend < 60)).any()
                condition2 = (diffWeekWeekend >= 60).any()
                condition3 = (diffWeekWeekend < -60).any()
            else:
                condition1 = -60 < diffWeekWeekend < 60
                condition2 = diffWeekWeekend >= 60
                condition3 = diffWeekWeekend < -60

            if condition1:
                act1Sleep = ""
            elif condition2:
                act1Sleep = "dormir mais durante a semana"
            elif condition3:
                act1Sleep = "dormir mais nos fins de semana"

    mSleepActions = ""
    if len(act1Sleep) > 0:
        mSleepActions = f"{mSleepActions}{act1Sleep}"
    if len(act2Sleep) > 0:
        if len(mSleepActions) > 0:
            mSleepActions = f"{mSleepActions}\n{act2Sleep}"
        else:
            mSleepActions = f"{mSleepActions}{act2Sleep}"
    if len(act3Sleep) > 0:
        if len(mSleepActions) > 0:
            mSleepActions = f"{mSleepActions}\n{act3Sleep}"
        else:
            mSleepActions = f"{mSleepActions}{act3Sleep}"

    # Return the results
    return act1Sleep, mSleepActions, mSleepDispersionGreen, mSleepDispersionRed, act2Sleep, mSleepQualityRed, mSleepQualityGreen, act3Sleep, avgCurrSleepMins, mSleepTimeRed, mSleepTimeGreen, mSleepTime

def calculate_steps(df, mpartID, currDateWeek, currWeekendArr, currWeekArr, lastDateWeek):
    """
    Function to calculate average steps, total activity time, and suggest actions based on step counts.
    """
    # Select necessary data
    currWeekSteps = df[(df['ParticipantProgramId'] == mpartID) & (df['DataType'] == "I") & (df['Date'] >= currDateWeek[0]) & (df['Date'] <= currDateWeek[1])]['HabitsSteps']
    lastWeekSteps = df[(df['ParticipantProgramId'] == mpartID) & (df['DataType'] == "I") & (df['Date'] >= lastDateWeek[0]) & (df['Date'] <= lastDateWeek[1])]['HabitsSteps']
    
    # Calculate averages
    avgCurrSteps = currWeekSteps.mean(numeric_only=True, skipna=True)
    avgLastSteps = lastWeekSteps.mean(numeric_only=True, skipna=True)

    # Handling NaN for averages and diff calculations
    avgCurrSteps = 0 if np.isnan(avgCurrSteps) else avgCurrSteps
    avgLastSteps = 0 if np.isnan(avgLastSteps) else avgLastSteps

    # Format the current week's average steps
    mSteps = "{:.1f}k".format(avgCurrSteps / 1000)

    # Calculate the difference in average steps
    diffSteps = avgCurrSteps - avgLastSteps

    # Conditionally format the difference based on data availability
    mStepsLastGreen = "+{:.1f}k".format(diffSteps / 1000) if diffSteps > 0 else ""
    
    if lastWeekSteps.empty:
        mStepsLastRed = "No Data Available"
    else:
        mStepsLastRed = "-{:.1f}k".format(abs(diffSteps) / 1000) if diffSteps < 0 else ""

    # Calculate activity sums and handle NaN
    def calculate_total_activity(activity_df):
        sumActivity = activity_df.sum(numeric_only=True, skipna=True)
        totalActivity = sumActivity.values[0] if not sumActivity.empty else 0
        return totalActivity / 60  # Convert minutes to hours

    # Current and last week activities
    currVigorous = df[(df['ParticipantProgramId'] == mpartID) & (df['DataType'] == "I") & (df['Date'] >= currDateWeek[0]) & (df['Date'] <= currDateWeek[1])].filter(like="Vigorous")
    currModerate = df[(df['ParticipantProgramId'] == mpartID) & (df['DataType'] == "I") & (df['Date'] >= currDateWeek[0]) & (df['Date'] <= currDateWeek[1])].filter(like="Moderate")
    totalCurrActivity = calculate_total_activity(pd.concat([currVigorous, currModerate]))

    lastVigorous = df[(df['ParticipantProgramId'] == mpartID) & (df['DataType'] == "I") & (df['Date'] >= lastDateWeek[0]) & (df['Date'] <= lastDateWeek[1])].filter(like="Vigorous")
    lastModerate = df[(df['ParticipantProgramId'] == mpartID) & (df['DataType'] == "I") & (df['Date'] >= lastDateWeek[0]) & (df['Date'] <= lastDateWeek[1])].filter(like="Moderate")
    totalLastActivity = calculate_total_activity(pd.concat([lastVigorous, lastModerate]))

    # Format hours and minutes for current activity
    hours, mins = divmod(totalCurrActivity * 60, 60)
    mActivitytime = "{:02d}:{:02d}".format(int(hours), int(mins))

    # Calculate the difference in total activity
    diffActivity = totalCurrActivity - totalLastActivity
    absDiff = abs(diffActivity)

    # Ensure hours and minutes are integers and format them
    hours, mins = divmod(absDiff * 60, 60)
    activityDiffFormat = "+{:02d}:{:02d}" if diffActivity > 0 else "-{:02d}:{:02d}"
    mActivitytimeGreen = activityDiffFormat.format(int(hours), int(mins)) if diffActivity > 0 else ""

    if lastWeekSteps.empty:
        mActivitytimeRed = "No Data Available"
    else:
        mActivitytimeRed = activityDiffFormat.format(int(hours), int(mins)) if diffActivity < 0 else ""

    # Calculate the average steps taken during the current weekend
    avgWeekendSteps = 0
    countdayssteps = 0

    for i in currWeekendArr:
        steps = df[(df['ParticipantProgramId'] == mpartID) & (df['DataType'] == "I") & (df['Date'] == i)].filter(like="HabitsSteps")
        if not steps.empty:
            avgWeekendSteps += steps.sum(numeric_only=True)
            countdayssteps += 1

    if countdayssteps != 0:
        avgWeekendSteps /= countdayssteps

    # Calculate the average steps taken during the current week
    avgWeekSteps = 0
    countdayssteps = 0

    # Assuming 'currWeekArr' is a list of dates for the current week
    for i in currWeekArr:
        steps = df[(df['ParticipantProgramId'] == mpartID) & (df['DataType'] == "I") & (df['Date'] == i)].filter(like="HabitsSteps")
        if not steps.empty:
            avgWeekSteps += steps.sum(numeric_only=True)
            countdayssteps += 1
    
    if countdayssteps != 0:
        avgWeekSteps /= countdayssteps

    # Calculate the difference between average steps taken during the weekend and the week
    diffWeekendWeekSteps = avgWeekendSteps - avgWeekSteps

    # Calculate the difference in average steps between the current period and the last period
    diffPeriodsSteps = avgCurrSteps - avgLastSteps

    action1Steps = ""

    if avgWeekendSteps.isna().any() or avgWeekSteps.isna().any():
        action1Steps = "melhorar a inserção de dados"
    else:
        if np.isnan(avgCurrSteps):
            action1Steps = "melhorar a recolha de dados diária"
        else:
            if round(avgCurrSteps / 1000, 1) < 10:
                if not np.isnan(avgCurrSteps):
                    if ((-3000 <= diffWeekendWeekSteps) & (diffWeekendWeekSteps <= 3000)).any():
                        if diffPeriodsSteps < 0: 
                            action1Steps = "aumentar número de passos"
                        else:
                            action1Steps = "continuar a aumentar passos"
                    elif (diffWeekendWeekSteps < -3000).any():
                        if diffPeriodsSteps < 0: 
                            action1Steps = "aumentar passos nos fins de semana"
                        else:
                            action1Steps = "continuar a aumentar passos"
                    else:
                        if diffPeriodsSteps < 0: 
                            action1Steps = "aumentar passos nos dias úteis"
                        else:
                            action1Steps = "continuar a aumentar passos"
                else:
                    action1Steps = "continuar a aumentar passos"
            else:
                action1Steps = "manter ≥10k passos por dia"

    diffActivity = (totalCurrActivity - totalLastActivity) * 60
    totalCurrActivityMins = totalCurrActivity * 60

    if (totalCurrActivityMins < 240).any():
        if (diffActivity < 0).any():  # use .all() here if all elements should be less than 0
            action2Steps = "aumentar atividade moderada e vigorosa"
        else:
            action2Steps = "continuar a aumentar atividade moderada e vigorosa"
    else:
        action2Steps = "manter ≥04:00 de atividade moderada e vigorosa por semana"

    mActivitytimeActions = ""
    if len(action1Steps) > 0:
        mActivitytimeActions += action1Steps
    if len(action2Steps) > 0:
        if len(mActivitytimeActions) > 0:
            mActivitytimeActions += "\n" + action2Steps
        else:
            mActivitytimeActions += action2Steps

    return mActivitytimeActions, mActivitytime, mSteps, mActivitytimeRed, mActivitytimeGreen, mStepsLastRed, mStepsLastGreen

def process_nutrition_data(df, mpartID, currDateWeek, lastDateWeek):
    """
    This function processes nutrition data and calculates average scores and sums for different nutrition elements.
    """
    # Selecting Nutrition Data for the Current Week
    currWeekNutri = df[(df['ParticipantProgramId'] == mpartID) & (df['DataType'] == "G") & (df['Date'] >= currDateWeek[0]) & (df['Date'] <= currDateWeek[1])].filter(like="MyHabitsNutritionEvolution")
    
    # Selecting Nutrition Data for the Last Week
    lastWeekNutri = df[(df['ParticipantProgramId'] == mpartID) & (df['DataType'] == "G") & (df['Date'] >= lastDateWeek[0]) & (df['Date'] <= lastDateWeek[1])].filter(like="MyHabitsNutritionEvolution")
    
    # Calculating the Average Nutrition Score for the Current Week
    avgCurrNutri = currWeekNutri.mean(numeric_only=True)

    # Calculating the Average Nutrition Score for the Last Week
    avgLastNutri = lastWeekNutri.mean(numeric_only=True)

    # Formatting the Current Week's Average Nutrition Score
    # Here, we loop through the index and format each value
    mNutritionscore = {index: "{:.0f}".format(val) for index, val in avgCurrNutri.items()}
    
    # Calculating the difference in Average Nutrition Scores
    diffNutri = avgCurrNutri - avgLastNutri

    # Initialize dictionaries to store the results
    mNutritionscoreRed = {}
    mNutritionscoreGreen = {}
    
    # Iterate through each item in the diffNutri Series
    for index, diff in diffNutri.items():
        # Checking if the difference is NA
        if pd.isna(diff):
            mNutritionscoreRed[index] = "No Data Available"
            mNutritionscoreGreen[index] = "No Data Available"
        else:
            # If the difference is positive or zero
            if diff >= 0:
                mNutritionscoreRed[index] = " "
                mNutritionscoreGreen[index] = "+{:.0f}".format(diff)
            # If the difference is negative
            else:
                mNutritionscoreRed[index] = "{:.0f}".format(diff)
                mNutritionscoreGreen[index] = " "

    # Selecting Daily Meals Data for the Current Week
    currWeekMeals = df[(df['ParticipantProgramId'] == mpartID) & (df['DataType'] == "I") & (df['Date'] >= currDateWeek[0]) & (df['Date'] <= currDateWeek[1])].filter(like="NutritionDailyMeals")
    
    # Selecting Daily Meals Data for the Last Week
    lastWeekMeals = df[(df['ParticipantProgramId'] == mpartID) & (df['DataType'] == "I") & (df['Date'] >= lastDateWeek[0]) & (df['Date'] <= lastDateWeek[1])].filter(like="NutritionDailyMeals")
    
    # Calculating the Average Daily Meals for the Current Week
    avgCurrMeals = currWeekMeals.mean(numeric_only=True)
    
    # Calculating the Average Daily Meals for the Last Week
    avgLastMeals = lastWeekMeals.mean(numeric_only=True)
    
    # Initialize dictionary to store the results
    mDailymeals = {}
    
    # Iterate through each item in the avgCurrMeals Series
    for index, avg in avgCurrMeals.items():
        # Checking if the average is NA
        if pd.isna(avg):
            mDailymeals[index] = "-"
        else:
            mDailymeals[index] = "{:.1f}".format(avg)

    # Initialize dictionaries to store the results
    mDailymealsGreen = {}
    mDailymealsRed = {}

    # Iterate through each pair of values in the avgCurrMeals and avgLastMeals Series
    for index in avgCurrMeals.index:
        avgCurr = avgCurrMeals[index]
        avgLast = avgLastMeals[index]
        
        # Convert values to numeric, errors='coerce' will convert non-numeric values to NaN
        avgCurr = pd.to_numeric(avgCurr, errors='coerce')
        avgLast = pd.to_numeric(avgLast, errors='coerce')
        
        # Checking if either average is NA or 0
        if pd.isna(avgCurr) or pd.isna(avgLast) or avgLast == 0:
            mDailymealsGreen[index] = " "
            mDailymealsRed[index] = "-"
        else:
            # Calculating the difference in Average Daily Meals
            diffMeals = avgCurr - avgLast
            
            # If the difference is positive or zero
            if diffMeals >= 0:
                mDailymealsGreen[index] = "+{:.1f}".format(diffMeals)
                mDailymealsRed[index] = " "
            # If the difference is negative
            else:
                mDailymealsGreen[index] = " "
                mDailymealsRed[index] = "{:.1f}".format(diffMeals)
    
    # Define the vectors
    vecNutritionElements = ["NutritionGrainsCerealsHighSugar","NutritionGrainsCerealsWholeGrain","NutritionGrainsCerealsHighFiber","NutritionGrainsCerealsMuesliGranola","NutritionGrainsBreadBrownWholeGrain","NutritionGrainsBreadWhite","NutritionGrainsRice_PastaWholeGrain","NutritionGrainsRice_PastaRegularGrain","NutritionFruitsBerries","NutritionFruitsNutsSeedsWithSalt","NutritionFruitsNutsSeedsWithoutSalt","NutritionFruitsOtherFruits","NutritionDairyYogurtWithSugar","NutritionDairyYogurtWithoutSugar","NutritionDairyYogurtSkir_Kefir_Lactobacilli","NutritionDairyYogurtGreek","NutritionDairyCheeseFullFat","NutritionDairyCheeseLowFat","NutritionDairyMilkWhole","NutritionDairyLowFatFatFree","NutritionDairyNonDairyMilk","NutritionMeatEggsBoiled","NutritionMeatEggsFriedScrambled","NutritionMeatRed","NutritionMeatWhite","NutritionMeatFish","NutritionVegetablesBeansLegumes","NutritionVegetablesPotatoesSweet","NutritionVegetablesOtherVegetables","NutritionVegetablesSoups","NutritionFatButterMargarine","NutritionFatLightButterVegetableOil","NutritionFatOliveOilOtherHealthyFats","NutritionFatCookingOil","NutritionSweetsHealthySweets","NutritionSweetsUnhealthySweets","NutritionBeverageWaterNaturalSparkling","NutritionBeverageWaterInfused","NutritionBeverageNaturalFruitJuicesSmoothies","NutritionBeverageVegetableJuicesSmoothies","NutritionBeverageSoftDrinksWithSugar","NutritionBeverageSoftDrinksWithoutSugar","NutritionBeverageAlcoholWineBeer","NutritionBeverageAlcoholSpirits","NutritionBeverageCoffeeDerivativesWithSugar","NutritionBeverageCoffeeDerivativesWithoutSugar","NutritionBeverageTeaWithSugar","NutritionBeverageTeaWithoutSugar","NutritionOtherArtificialProteinSourcesSupplements","NutritionOtherSaucesCondiments","NutritionOtherFastFood"]
    vecNutritionLabels = ["Cereais c/açúcar","Cereais int.","Cereais c/fibra","Muesli/Granola","Pão não-branco","Pão branco","Arroz/Massa int.","Arroz/Massa norm.","Fr. silvestres","Fr. secos/sem. c/sal"," Fr. secos/sem. s/sal","Fruta","Iog. c/açúcar","Iog. s/açúcar","Skir/Kefir/Lbac.","Iog. grego","Queijo gordo","Queijo magro","Leite gordo","Leite não-gordo","Bebida veg.","Ovo cozido","Ovo mexido","Carne verm. ","Carne branca","Peixe","Feijão/Legumin.","Batata doce","Vegetais","Sopa","Manteiga/Marg.","Mant.Mag./Ól.Veg.","Azeite/Gord.saud.","Óleo alimentar","Doces saud.","Doces não-saud.","Água","Infusão","Sumo/smoothie frut."," Sumo/smoothie veg.","Refrig. c/açúcar","Refrig. s/açúcar","Vinho/Cerveja","Beb. espirituosas","Café c/açúcar","Café s/açúcar","Chá c/açúcar","Chá s/açúcar","Prot.art./Supl.","Molhos/condim.","Fast Food"]
    vecNutritionWeights = [-1,1,1,1,1,0,1,0,1,-1,1,1,-1,1,1,-1,-1,1,-1,1,1,1,1,0,1,1,1,1,1,1,-1,-1,1,-1,0,-1,1,1,1,1,-1,0,-1,-1,-1,1,-1,1,1,-1,-1]

    # Initialize an empty list
    vecSumNutritionElements = []

    # Iterate over each element in vecNutritionElements
    for element in vecNutritionElements:
        # Select the corresponding data from the DataFrame
        Nutri = df[(df['ParticipantProgramId'] == mpartID) & (df['DataType'] == "I") & (df['Date'] >= currDateWeek[0]) & (df['Date'] <= currDateWeek[1])].filter(like=element)
        # Calculate the sum and append it to vecSumNutritionElements
        NutriSum = Nutri.sum().sum()
        vecSumNutritionElements.append(int(NutriSum))

    # Create a DataFrame
    dfNutri = pd.DataFrame({
        "names": vecNutritionLabels,
        "sumVals": vecSumNutritionElements,
        "othernames": vecNutritionElements,
        "weights": vecNutritionWeights
    })

    # Select rows where sumVals > 0 and weights > 0
    dfNutri_Positive = dfNutri[(dfNutri['sumVals'] > 0) & (dfNutri['weights'] > 0)][['names', 'sumVals']]
    # Select rows where sumVals > 0 and weights < 0
    dfNutri_Negative = dfNutri[(dfNutri['sumVals'] > 0) & (dfNutri['weights'] < 0)][['names', 'sumVals']]

    # Sort dfNutri_Positive in descending order by 'sumVals'
    dfNutriPosOrder = dfNutri_Positive.sort_values(by='sumVals', ascending=False)
    # Sort dfNutri_Negative in descending order by 'sumVals'
    dfNutriNegOrder = dfNutri_Negative.sort_values(by='sumVals', ascending=False)

    # Initialize the variables - NEGATIVE
    mRedNutrition1 = mRedLabelNutrition1 = " "
    mRedNutrition2 = mRedLabelNutrition2 = " "
    mRedNutrition3 = mRedLabelNutrition3 = " "
    mRedNutrition4 = mRedLabelNutrition4 = " "
    mRedNutrition5 = mRedLabelNutrition5 = " "

    # Assign the values and labels if they are not NA
    if dfNutriNegOrder.shape[0] > 0 and dfNutriNegOrder.iloc[0, 0] is not pd.NA:
        mRedNutrition1 = int(dfNutriNegOrder.iloc[0, 1])  
        mRedLabelNutrition1 = dfNutriNegOrder.iloc[0, 0]
    if dfNutriNegOrder.shape[0] > 1 and dfNutriNegOrder.iloc[1, 0] is not pd.NA:
        mRedNutrition2 = int(dfNutriNegOrder.iloc[1, 1])  
        mRedLabelNutrition2 = dfNutriNegOrder.iloc[1, 0]
    if dfNutriNegOrder.shape[0] > 2 and dfNutriNegOrder.iloc[2, 0] is not pd.NA:
        mRedNutrition3 = int(dfNutriNegOrder.iloc[2, 1])  
        mRedLabelNutrition3 = dfNutriNegOrder.iloc[2, 0]
    if dfNutriNegOrder.shape[0] > 3 and dfNutriNegOrder.iloc[3, 0] is not pd.NA:
        mRedNutrition4 = int(dfNutriNegOrder.iloc[3, 1])  
        mRedLabelNutrition4 = dfNutriNegOrder.iloc[3, 0]
    if dfNutriNegOrder.shape[0] > 4 and dfNutriNegOrder.iloc[4, 0] is not pd.NA:
        mRedNutrition5 = int(dfNutriNegOrder.iloc[4, 1])  
        mRedLabelNutrition5 = dfNutriNegOrder.iloc[4, 0]

    # Initialize the variables - POSITIVE
    mGreenNutrition1 = mNutrition1 = " "
    mGreenNutrition2 = mNutrition2 = " "
    mGreenNutrition3 = mNutrition3 = " "
    mGreenNutrition4 = mNutrition4 = " "
    mGreenNutrition5 = mNutrition5 = " "

    # Assign the values and labels if they are not NA
    if dfNutriPosOrder.shape[0] > 0 and dfNutriPosOrder.iloc[0, 0] is not pd.NA:
        mGreenNutrition1 = int(dfNutriPosOrder.iloc[0, 1])  
        mNutrition1 = dfNutriPosOrder.iloc[0, 0]
    if dfNutriPosOrder.shape[0] > 1 and dfNutriPosOrder.iloc[1, 0] is not pd.NA:
        mGreenNutrition2 = int(dfNutriPosOrder.iloc[1, 1])  
        mNutrition2 = dfNutriPosOrder.iloc[1, 0]
    if dfNutriPosOrder.shape[0] > 2 and dfNutriPosOrder.iloc[2, 0] is not pd.NA:
        mGreenNutrition3 = int(dfNutriPosOrder.iloc[2, 1])  
        mNutrition3 = dfNutriPosOrder.iloc[2, 0]
    if dfNutriPosOrder.shape[0] > 3 and dfNutriPosOrder.iloc[3, 0] is not pd.NA:
        mGreenNutrition4 = int(dfNutriPosOrder.iloc[3, 1])  
        mNutrition4 = dfNutriPosOrder.iloc[3, 0]
    if dfNutriPosOrder.shape[0] > 4 and dfNutriPosOrder.iloc[4, 0] is not pd.NA:
        mGreenNutrition5 = int(dfNutriPosOrder.iloc[4, 1])  
        mNutrition5 = dfNutriPosOrder.iloc[4, 0]
    
    # Set the value of fracValueComparison
    fracValueComparison = 0.8
    
    return mNutritionscore, mNutritionscoreRed, mNutritionscoreGreen, avgCurrMeals, mDailymeals, mDailymealsGreen, mDailymealsRed, mRedNutrition1, mRedNutrition2, mRedNutrition3, mRedNutrition4, mRedNutrition5, mRedLabelNutrition1, mRedLabelNutrition2, mRedLabelNutrition3, mRedLabelNutrition4, mRedLabelNutrition5, mNutrition1, mNutrition2, mNutrition3, mNutrition4, mNutrition5, mGreenNutrition1, mGreenNutrition2, mGreenNutrition3, mGreenNutrition4, mGreenNutrition5, fracValueComparison

def calculate_nutrition_actions(df, mpartID, currWeekendArr, currWeekArr):
    """
    Calculate the average nutritional score for the weekend and the week using the data in df DataFrame.
    Then determine an action based on the difference between the week and weekend averages.
    """
    avgWeekendNutri = 0
    countNutriDays = 0
    for date in currWeekendArr:
        data_row = df.loc[(df['ParticipantProgramId'] == mpartID) & (df['DataType'] == "G") & (df['Date'] == date)]
        for column in data_row.columns:
            if "MyHabitsNutritionEvolution" in column and pd.notna(data_row[column].values[0]):
                avgWeekendNutri += data_row[column].values[0]
                countNutriDays += 1
    
    if countNutriDays > 0:
        avgWeekendNutri = avgWeekendNutri / countNutriDays
    else:
        avgWeekendNutri = None  # Handle no data scenario

    avgWeekNutri = 0
    countNutriDays = 0
    for date in currWeekArr:
        data_row = df.loc[(df['ParticipantProgramId'] == mpartID) & (df['DataType'] == "G") & (df['Date'] == date)]
        for column in data_row.columns:
            if "MyHabitsNutritionEvolution" in column and pd.notna(data_row[column].values[0]):
                avgWeekNutri += data_row[column].values[0]
                countNutriDays += 1

    if countNutriDays > 0:
        avgWeekNutri = avgWeekNutri / countNutriDays
    else:
        avgWeekNutri = None  # Handle no data scenario

    if avgWeekendNutri is None or avgWeekNutri is None:
        action1Nutri = "melhorar a inserção de dados de nutrição"
    else:
        diffWeekWeekendNutri = avgWeekNutri - avgWeekendNutri
        action1Nutri = ""
        if diffWeekWeekendNutri >= 15:
            action1Nutri = "melhorar as escolhas nutricionais ao fim de semana"
        elif diffWeekWeekendNutri <= -10:
            action1Nutri = "melhorar as escolhas nutricionais nos dias úteis"

    return action1Nutri

def calculate_meal_actions(avgCurrMeals):
    """
    Determine the recommended action based on the average number of current meals.
    """
    action2Nutri = ""

    # Ensure avgCurrMeals is treated as a Series and handle conditions
    if isinstance(avgCurrMeals, pd.Series):
        if pd.isna(avgCurrMeals).any():
            action2Nutri = ""  # No action if there are NaN values
        elif (avgCurrMeals < 1).any():
            action2Nutri = "fazer ≥4 refeições por dia"  # Action if any value is less than 1
        else:
            action2Nutri = ""  # No action if all values are 1 or more
    else:
        # Handle the case where avgCurrMeals might not be a Series
        if pd.isna(avgCurrMeals):
            action2Nutri = ""
        elif avgCurrMeals < 1:
            action2Nutri = "fazer ≥4 refeições por dia"
        else:
            action2Nutri = ""

    return action2Nutri

def calculate_fruit_actions(df, mpartID, currDateWeek, fracValueComparison):
    """
    Calculate the recommended action based on the consumption of fruits.
    """
    action3Nutri = ""
    
    # Filter the dataframe for the specific participant and date range
    Fruit = df.loc[
        (df['ParticipantProgramId'] == mpartID) & 
        (df['DataType'] == "G") & 
        (df['Date'] >= currDateWeek[0]) & 
        (df['Date'] <= currDateWeek[1]), 
        df.columns[df.columns.str.startswith("NutritionResultFreshFruit")]
    ]

    goodFruit = 0
    badFruit = 0

    # Iterate through the rows in the filtered dataframe
    for i in range(len(Fruit)):
        fruit_value = Fruit.iloc[i, 0]
        if pd.isna(fruit_value):  # Check for NaN values
            continue  # Skip this iteration if the value is NaN
        
        if int(fruit_value) - 9000000 == 0:
            g = 0
            b = 0
        else:
            # Extract parts of the number safely
            fruit_str = str(int(fruit_value))
            g = int(fruit_str[1:4]) if len(fruit_str) > 3 else 0
            b = int(fruit_str[4:7]) if len(fruit_str) > 6 else 0
        
        goodFruit += g
        badFruit += b

    # Calculate the fraction if possible
    if goodFruit == 0 and badFruit == 0:
        frac = -1
    else:
        frac = goodFruit / badFruit if badFruit != 0 else float('inf')  # Avoid division by zero

    # Determine the action based on the fraction
    if 0 <= frac < fracValueComparison:
        action3Nutri = "comer fruta ≥4 vezes por dia"

    return action3Nutri

def calculate_sugar_actions(df, mpartID, currDateWeek, fracValueComparison):
    action4Nutri = ""

    Sugar = df.loc[
        (df['ParticipantProgramId'] == mpartID) & 
        (df['DataType'] == "G") & 
        (df['Date'] >= currDateWeek[0]) & 
        (df['Date'] <= currDateWeek[1]), 
        df.columns[df.columns.str.startswith("NutritionResultSugarLevels")]
    ]

    goodSugarVal = 0
    badSugarVal = 0

    for i in range(len(Sugar)):
        sugar_value = Sugar.iloc[i, 0]
        if pd.isna(sugar_value):
            continue

        sugar_value = int(sugar_value)
        if sugar_value - 9000000 == 0:
            g = 0
            b = 0
        else:
            sugar_str = str(sugar_value)
            g = int(sugar_str[1:4]) if len(sugar_str) > 3 else 0
            b = int(sugar_str[4:7]) if len(sugar_str) > 6 else 0

        goodSugarVal += g
        badSugarVal += b

    if goodSugarVal == 0 and badSugarVal == 0:
        return "No action necessary due to zero consumption."

    frac = goodSugarVal / badSugarVal if badSugarVal != 0 else float('inf')

    if frac < fracValueComparison or (frac == float('inf') and goodSugarVal > 0):
        action4Nutri = "substituir doces por frutos silvestres, frutos secos ou sementes"

    return action4Nutri

def calculate_whole_grain_actions(df, mpartID, currDateWeek, fracValueComparison):
    """
    Calculate the recommended action based on the consumption of whole grains.
    """
    action5Nutri = ""

    # Filter the DataFrame for the specific participant and date range
    WholeGrain = df.loc[
        (df['ParticipantProgramId'] == mpartID) & 
        (df['DataType'] == "G") & 
        (df['Date'] >= currDateWeek[0]) & 
        (df['Date'] <= currDateWeek[1]), 
        df.columns[df.columns.str.startswith("NutritionResultWholeGrainFood")]
    ]

    goodWholeGrain = 0
    badWholeGrain = 0

    # Process each row in the filtered DataFrame
    for i in range(len(WholeGrain)):
        # Ensure the value is not NaN before processing
        if pd.isna(WholeGrain.iloc[i, 0]):
            continue  # Skip this iteration if the value is NaN

        # Safe conversion to integer
        whole_grain_value = int(WholeGrain.iloc[i, 0])

        if whole_grain_value - 9000000 == 0:
            g = 0
            b = 0
        else:
            whole_grain_str = str(whole_grain_value)
            g = int(whole_grain_str[1:4]) if len(whole_grain_str) > 3 else 0
            b = int(whole_grain_str[4:7]) if len(whole_grain_str) > 6 else 0

        goodWholeGrain += g
        badWholeGrain += b

    # Calculate the fraction of good to bad whole grain
    if goodWholeGrain == 0 and badWholeGrain == 0:
        frac = -1  # Indicates no consumption
    else:
        frac = goodWholeGrain / badWholeGrain if badWholeGrain != 0 else float('inf')  # Handle division by zero

    # Determine the nutritional action based on the fraction
    if 0 <= frac < fracValueComparison:
        action5Nutri = "preferir cereais, pão, massa e arroz integrais"

    return action5Nutri

def calculate_protein_actions(df, mpartID, currDateWeek, fracValueComparison):
    """
    Calculate the recommended action based on the consumption of proteins.
    """
    action6Nutri = ""

    # Filter the DataFrame for the specific participant and date range
    Protein = df.loc[
        (df['ParticipantProgramId'] == mpartID) & 
        (df['DataType'] == "G") & 
        (df['Date'] >= currDateWeek[0]) & 
        (df['Date'] <= currDateWeek[1]), 
        df.columns[df.columns.str.startswith("NutritionResultGoodProteinSources")]
    ]

    goodProtein = 0
    badProtein = 0

    # Process each row in the filtered DataFrame
    for i in range(len(Protein)):
        protein_value = Protein.iloc[i, 0]

        # Ensure the value is not NaN before processing
        if pd.isna(protein_value):
            continue  # Skip this iteration if the value is NaN

        # Safe conversion to integer
        protein_value = int(protein_value)

        if protein_value - 9000000 == 0:
            g = 0
            b = 0
        else:
            protein_str = str(protein_value)
            g = int(protein_str[1:4]) if len(protein_str) > 3 else 0
            b = int(protein_str[4:7]) if len(protein_str) > 6 else 0

        goodProtein += g
        badProtein += b

    # Calculate the fraction of good to bad protein
    if goodProtein == 0 and badProtein == 0:
        frac = -1  # Indicates no consumption
    else:
        frac = goodProtein / badProtein if badProtein != 0 else float('inf')  # Handle division by zero

    # Determine the nutritional action based on the fraction
    if 0 <= frac < fracValueComparison:
        action6Nutri = "preferir peixe a carnes vermelhas"

    return action6Nutri

def calculate_water_actions(df, mpartID, currDateWeek, fracValueComparison):
    """
    Calculate the recommended action based on the consumption of water.
    """
    action7Nutri = ""

    # Filter the DataFrame for the specific participant and date range
    Water = df.loc[
        (df['ParticipantProgramId'] == mpartID) & 
        (df['DataType'] == "G") & 
        (df['Date'] >= currDateWeek[0]) & 
        (df['Date'] <= currDateWeek[1]), 
        df.columns[df.columns.str.startswith("NutritionResultWaterLevels")]
    ]

    goodWater = 0
    badWater = 0

    # Process each row in the filtered DataFrame
    for i in range(len(Water)):
        water_value = Water.iloc[i, 0]

        # Ensure the value is not NaN before processing
        if pd.isna(water_value):
            continue  # Skip this iteration if the value is NaN

        # Safe conversion to integer
        water_value = int(water_value)

        if water_value - 9000000 == 0:
            g = 0
            b = 0
        else:
            water_str = str(water_value)
            g = int(water_str[1:4]) if len(water_str) > 3 else 0
            b = int(water_str[4:7]) if len(water_str) > 6 else 0

        goodWater += g
        badWater += b

    # Calculate the fraction of good to bad water
    if goodWater == 0 and badWater == 0:
        frac = -1  # Indicates no consumption
    else:
        frac = goodWater / badWater if badWater != 0 else float('inf')  # Handle division by zero

    # Determine the nutritional action based on the fraction
    if 0 <= frac < fracValueComparison:
        action7Nutri = "preferir beber água natural ou bebidas sem açúcar"

    return action7Nutri

def calculate_vegetable_actions(df, mpartID, currDateWeek):
    """
    Calculate the recommended action based on the consumption of vegetables.
    """
 
    action8Nutri = ""
 
    veg_columns = [col for col in df.columns if col.startswith("NutritionVegetables")]
   
    df_veg = df.copy()
    df_veg[veg_columns] = df_veg[veg_columns].fillna(0)
 
    veg1 = df_veg.loc[(df_veg['ParticipantProgramId']==mpartID) &
                (df_veg['DataType']=="I") &
                (df_veg['Date'] >= currDateWeek[0]) &
                (df_veg['Date'] <= currDateWeek[1]),
                df_veg.columns.str.startswith("NutritionVegetablesBeansLegumes")]
 
    veg2 = df_veg.loc[(df_veg['ParticipantProgramId']==mpartID) &
                (df_veg['DataType']=="I") &
                (df_veg['Date'] >= currDateWeek[0]) &
                (df_veg['Date'] <= currDateWeek[1]),
                df_veg.columns.str.startswith("NutritionVegetablesOtherVegetables")]
 
    veg3 = df_veg.loc[(df_veg['ParticipantProgramId']==mpartID) &
                (df_veg['DataType']=="I") &
                (df_veg['Date'] >= currDateWeek[0]) &
                (df_veg['Date'] <= currDateWeek[1]),
                df_veg.columns.str.startswith("NutritionVegetablesSoups")]
   
    totalVeg = veg1.values.sum() + veg2.values.sum() + veg3.values.sum()
    totalVeg = totalVeg / len(veg1)
 
    if totalVeg < 2 and totalVeg != 0:
        action8Nutri = "comer mais legumes e vegetais"
 
    return action8Nutri

def consolidate_nutrition_actions(action1Nutri, action2Nutri, action3Nutri, action4Nutri, action5Nutri, action6Nutri, action7Nutri, action8Nutri):
    """
    Consolidates multiple nutrition action strings into a single string.
    
    This function takes up to eight nutrition action strings as input. It filters out any empty strings, and then joins the remaining 
    action strings together, separated by newline characters. The resulting string is returned.
    """
    
    mActionsNutrition = []
    
    for action in [action1Nutri, action2Nutri, action3Nutri, action4Nutri, action5Nutri, action6Nutri, action7Nutri, action8Nutri]:
        if len(action) > 0:
            mActionsNutrition.append(action)

    mActionsNutrition = "\n".join(mActionsNutrition)

    return mActionsNutrition

from pptx import Presentation

def list_placeholders_in_ppt(presentation_path):
    # Load the presentation
    prs = Presentation(presentation_path)
    
    # Iterate over all slides
    for slide_number, slide in enumerate(prs.slides, start=1):
        # Iterate over all placeholders in the slide
        for placeholder in slide.placeholders:
            # Fetching the placeholder format and type
            ph_format = placeholder.placeholder_format
            print(f"Placeholder index: {placeholder.placeholder_format.idx}, name: {placeholder.name}")

def main():
    """
    Main function to orchestrate the entire process.
    """
    file_path = 'FUEL_onemonth_report.csv'
    
    # Load the data
    df = load_csv_to_dataframe(file_path)
    # Process 'Date' column
    df_sorted = process_date_column(df) 

    #Calculate week dates
    curr_date_week, last_date_week, curr_weekend_arr, curr_week_arr = calculate_week_dates(df_sorted)
    print("Current Week Dates:", curr_date_week)
    print("Last Week Dates:", last_date_week)
    print("Current Weekend Dates:", curr_weekend_arr)
    print("Current Weekday Dates:", curr_week_arr)

    # Preprocess participant data
    df_preprocess = pre_process_dataframe(df_sorted)

    # Analyze participant data and generate a PowerPoint report for each participant
    mpart_id = analyze_participant_data(df_preprocess, curr_date_week, curr_weekend_arr, curr_week_arr, last_date_week)

    list_placeholders_in_ppt("weekReport.pptx")

if __name__ == "__main__":
    main()




